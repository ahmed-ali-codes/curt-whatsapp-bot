import fitz  # PyMuPDF for PDF
import google.generativeai as genai
import os
from dotenv import load_dotenv
from docx import Document  # For Word files
from pptx import Presentation  # For PowerPoint files

# -----------------------------
# Load Environment Variables
# -----------------------------
load_dotenv()

# -----------------------------
# Gemini API Key (loaded from .env)
# -----------------------------
GEMINI_API_KEY = os.getenv("GEMINI_API_KEY")
if not GEMINI_API_KEY:
    raise ValueError("❌ GEMINI_API_KEY not found! Please set it in your .env file.")
genai.configure(api_key=GEMINI_API_KEY)

# Load Gemini Model
model = genai.GenerativeModel("gemini-2.5-flash")

# -----------------------------
# Extract Text from PDF
# -----------------------------
def extract_text_from_pdf(pdf_path):
    text = ""
    try:
        with fitz.open(pdf_path) as doc:
            for page in doc:
                text += page.get_text()
    except Exception as e:
        print(f"❌ Error reading PDF: {e}")
        return None
    return text.strip()

# -----------------------------
# Extract Text from Word (.docx)
# -----------------------------
def extract_text_from_word(docx_path):
    text = ""
    try:
        doc = Document(docx_path)
        for para in doc.paragraphs:
            text += para.text + "\n"
    except Exception as e:
        print(f"❌ Error reading Word file: {e}")
        return None
    return text.strip()

# -----------------------------
# Extract Text from PowerPoint (.pptx)
# -----------------------------
def extract_text_from_pptx(pptx_path):
    text = ""
    try:
        prs = Presentation(pptx_path)
        for slide in prs.slides:
            for shape in slide.shapes:
                if hasattr(shape, "text"):
                    text += shape.text + "\n"
    except Exception as e:
        print(f"❌ Error reading PowerPoint file: {e}")
        return None
    return text.strip()

# -----------------------------
# Function: Extract Text Automatically
# -----------------------------
def extract_text(file_path):
    ext = os.path.splitext(file_path)[1].lower()
    if ext == ".pdf":
        return extract_text_from_pdf(file_path)
    elif ext == ".docx":
        return extract_text_from_word(file_path)
    elif ext == ".pptx":
        return extract_text_from_pptx(file_path)
    else:
        print("❌ Unsupported file format! Please upload a PDF, Word (.docx), or PowerPoint (.pptx) file.")
        return None

# -----------------------------
# Generate with Gemini
# -----------------------------
def generate_with_gemini(task, text):
    if task == "1":
        prompt = f"Generate 10 short questions and answers from this study material:\n\n{text}"
    elif task == "2":
        prompt = f"Create 10 multiple-choice questions (each with 4 options and one correct answer) from this study material:\n\n{text}"
    elif task == "3":
        prompt = f"Summarize this study material into concise key points:\n\n{text}"
    else:
        return "❌ Invalid task selection."

    try:
        response = model.generate_content(prompt)
        return response.text.strip()
    except Exception as e:
        return f"❌ Error generating content with Gemini: {e}"

# -----------------------------
# Save Output to File
# -----------------------------
def save_output(task, content):
    file_names = {
        "1": "output_shortqa.txt",
        "2": "output_mcqs.txt",
        "3": "output_summary.txt"
    }
    filename = file_names.get(task, "output.txt")
    with open(filename, "w", encoding="utf-8") as f:
        f.write(content)
    print(f"\n💾 Output saved as: {filename}")

# -----------------------------
# Main Loop
# -----------------------------
def main():
    print("=== 📘 Study Material Assistant (PDF, Word, PPT → Gemini AI) ===\n")
    while True:
        print("\nSelect a task:")
        print("1️⃣  Generate Short Questions & Answers")
        print("2️⃣  Generate MCQs (Multiple Choice Questions)")
        print("3️⃣  Summarize the File")
        print("4️⃣  Exit")

        task = input("\n👉 Enter your choice (1/2/3/4): ").strip()
        if task == "4" or task.lower() == "exit":
            print("\n👋 Exiting... Goodbye!\n")
            break

        file_path = input("\n📂 Enter the file path (.pdf / .docx / .pptx): ").strip()
        if not os.path.exists(file_path):
            print("❌ File not found! Please try again.")
            continue

        print("\n⏳ Extracting text from file...")
        text = extract_text(file_path)
        if not text:
            print("❌ Could not extract text. Please check file content.")
            continue

        print("🤖 Generating response from Gemini... Please wait...")
        result = generate_with_gemini(task, text)

        print("\n✅ Result:\n")
        print(result)
        print("\n" + "="*70)

        save_output(task, result)

# -----------------------------
# Run the Program
# -----------------------------
if __name__ == "__main__":
    main()
